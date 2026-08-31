from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 16:9 와이드스크린 설정
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# 색상 팔레트 (화이트 / 모던 클린 테마)
BG_COLOR = RGBColor(255, 255, 255)
CARD_BG = RGBColor(248, 250, 252)
CARD_BORDER = RGBColor(226, 232, 240)
PRIMARY_BLUE = RGBColor(37, 99, 235)
ACCENT_BLUE = RGBColor(239, 246, 255)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)
TEXT_SUB = RGBColor(71, 85, 105)
ACCENT_GREEN = RGBColor(5, 150, 105)
TABLE_HEADER_BG = RGBColor(241, 245, 249)

def apply_background(slide):
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_COLOR
    bg_shape.line.fill.background()

def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_header(slide, category, title):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_cat = tf.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = PRIMARY_BLUE
    
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK
    p_title.space_before = Pt(4)

# ==========================================
# [Slide 1] 표지
# ==========================================
s1 = prs.slides.add_slide(blank_layout)
apply_background(s1)

tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "AI / ML ENGINEER PORTFOLIO"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

p = tf1.add_paragraph()
p.text = "정량 지표 중심의 모델링과\n실서비스 서빙 최적화를 완성하는 엔지니어"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = TEXT_DARK
p.space_before = Pt(12)

p = tf1.add_paragraph()
p.text = "조용준 | AI & ML Engineering"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_SUB
p.space_before = Pt(20)

p = tf1.add_paragraph()
p.text = "• Email: dydwns2663@naver.com\n• GitHub: https://github.com/Eden-Cho\n• Blog: https://velog.io/@goom95"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(10)

# ==========================================
# [Slide 2] 프로젝트 1 개요
# ==========================================
s2 = prs.slides.add_slide(blank_layout)
apply_background(s2)
add_header(s2, "Project 01 • Overview & Pipeline", "소상공인 AI 광고 상세페이지 생성 서비스")

create_card(s2, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.4))
tb = s2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📌 프로젝트 개요 및 핵심 기술"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

p = tf.add_paragraph()
p.text = "• 기간 / 역할: 2026.06 - 2026.07 (5인 팀 / AI 백엔드 서빙 & MLOps 담당)\n• 핵심 기술: Python, FastAPI, PyTorch, Langfuse, CLIP Score, BRISQUE, Docker, GCP\n• 주요 목표: 멀티모달 자동 품질 평가 파이프라인 구축 및 무거운 연산의 비동기 서빙 안정화"
p.font.size = Pt(11)
p.font.color.rgb = TEXT_SUB
p.space_before = Pt(4)

steps = [
    ("Step 1. 유저 요청", "Streamlit UI\n광고 파라미터 입력 및 생성"),
    ("Step 2. 서빙 API", "FastAPI 엔드포인트\n요청 검증 및 파이프라인 라우팅"),
    ("Step 3. 비동기 분리", "BackgroundTasks\nCLIP / BRISQUE 평가 비동기화"),
    ("Step 4. MLOps 관측", "Langfuse Tracing\n비용/지연 실시간 대시보드")
]
for i, (stitle, sdesc) in enumerate(steps):
    left = Inches(0.8 + i * 3.0)
    create_card(s2, left, Inches(3.5), Inches(2.73), Inches(3.2))
    tb_step = s2.shapes.add_textbox(left + Inches(0.2), Inches(3.8), Inches(2.33), Inches(2.6))
    tf_step = tb_step.text_frame
    tf_step.word_wrap = True
    p1 = tf_step.paragraphs[0]
    p1.text = stitle
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_BLUE
    p2 = tf_step.add_paragraph()
    p2.text = sdesc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SUB
    p2.space_before = Pt(8)

# ==========================================
# [Slide 3] 프로젝트 1 최적화 & MLOps (수치 비교 보강)
# ==========================================
s3 = prs.slides.add_slide(blank_layout)
apply_background(s3)
add_header(s3, "Project 01 • Deep Dive & Results", "서빙 최적화 및 정량 지표 개선 성과")

rows, cols = 4, 4
table_shape = s3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.6))
table = table_shape.table
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(3.1)
table.columns[2].width = Inches(3.3)
table.columns[3].width = Inches(3.53)

headers = ["핵심 영역", "Before (초기 베이스라인)", "기술적 의사결정 (Action)", "After (정량 개선 성과)"]
for col_idx, h_text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER_BG
    p = cell.text_frame.paragraphs[0]
    p.text = h_text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

table_data = [
    ("Cold Start", "첫 호출 시 가중치 로딩으로\n초기 20~30초 응답 지연", "Lifespan 이벤트 기반\n인메모리 Pre-loading 적용", "초기 로딩 0초 즉시 응답\n(Cold Start 완전 해소)"),
    ("서빙 지연", "동기식 화질 평가 연산으로\n평균 응답 3.8초 대기 발생", "FastAPI BackgroundTasks로\n품질 평가 비동기 분리", "체감 응답 1.5초 (약 60.5% 단축)\n동시 요청(20명) 성공률 100%"),
    ("화질 / 비용", "평가 모니터링 부재 및\n화질 지표 BRISQUE 38.3", "파이프라인 v2 고도화 &\nLangfuse @observe 트레이싱", "BRISQUE 31.4 (18.0% 화질 개선)\n1회당 $0.0082(11원) 비용 통제")
]

for row_idx, row_data in enumerate(table_data, start=1):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_SUB
        if col_idx == 3:
            p.font.bold = True
            p.font.color.rgb = ACCENT_GREEN

create_card(s3, Inches(0.8), Inches(5.7), Inches(11.73), Inches(1.1), bg_color=ACCENT_BLUE, border_color=PRIMARY_BLUE)
tb_sub = s3.shapes.add_textbox(Inches(1.1), Inches(5.8), Inches(11.1), Inches(0.8))
tf_sub = tb_sub.text_frame
tf_sub.word_wrap = True
p = tf_sub.paragraphs[0]
p.text = "📊 핵심 요약: 응답 대기 시간 60.5% 단축(3.8s ➔ 1.5s) 및 무참조 화질 지표 18.0% 개선(38.3 ➔ 31.4)을 정량 검증함"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

# ==========================================
# [Slide 4] 프로젝트 2 개요
# ==========================================
s4 = prs.slides.add_slide(blank_layout)
apply_background(s4)
add_header(s4, "Project 02 • Overview & Pipeline", "경구 알약 객체 탐지 모델 개발 및 아키텍처 비교")

create_card(s4, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.3))
tb4 = s4.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(1.1))
tf4 = tb4.text_frame
tf4.word_wrap = True
p = tf4.paragraphs[0]
p.text = "📌 프로젝트 개요 및 문제 정의"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

p = tf4.add_paragraph()
p.text = "• 기간 / 역할: 2026.03 - 2026.04 (5인 팀 / CV 모델링 & 데이터 파이프라인 담당)\n• 문제 상황: 73종 알약의 미세 각인, 소형 객체, 심각한 데이터 불균형으로 1-Stage 단독 모델 한계 봉착"
p.font.size = Pt(11)
p.font.color.rgb = TEXT_SUB
p.space_before = Pt(3)

p2_steps = [
    ("Stage 1. 위치 검출", "1-Class Detector (YOLO26m)\n• 73개 클래스 구분 없이 객체 탐지\n• BBox 검출 및 패치 자동 Crop"),
    ("Stage 2. 정밀 분류", "Classifier (EfficientNet-B3)\n• 73종 세부 품목 및 미세 각인 분류\n• WeightedRandomSampler 적용"),
    ("최종 성과 도출", "2-Stage 통합 하이브리드\n• 1-Stage 대비 정확도 14.5% 향상\n• Kaggle 0.987 최고 스코어 달성")
]
for i, (stitle, sdesc) in enumerate(p2_steps):
    left = Inches(0.8 + i * 4.0)
    create_card(s4, left, Inches(3.4), Inches(3.73), Inches(3.3))
    tb_step = s4.shapes.add_textbox(left + Inches(0.2), Inches(3.7), Inches(3.33), Inches(2.7))
    tf_step = tb_step.text_frame
    tf_step.word_wrap = True
    p1 = tf_step.paragraphs[0]
    p1.text = stitle
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_BLUE
    p2 = tf_step.add_paragraph()
    p2.text = sdesc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SUB
    p2.space_before = Pt(8)

# ==========================================
# [Slide 5] 프로젝트 2 모델 실험 및 튜닝 (수치 비교 보강)
# ==========================================
s5 = prs.slides.add_slide(blank_layout)
apply_background(s5)
add_header(s5, "Project 02 • Experiments & Tuning", "정량적 모델 최적화 및 72회 그리드 서치")

create_card(s5, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.9))
tb_l = s5.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
tf_l = tb_l.text_frame
tf_l.word_wrap = True
p = tf_l.paragraphs[0]
p.text = "🧪 불균형 교정 및 백본 튜닝"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

bullets_l = [
    "• 도메인 맞춤 증강 (Albumentations):\n  알약의 형상 및 미세 각인을 보존하는 증강 설계",
    "• WeightedRandomSampler 도입:\n  소수 품목 가중치를 보정하여 클래스 불균형 교정\n  ➔ 1-Stage 단독 대비 분류 정확도 14.5% 향상",
    "• Faster R-CNN 백본 튜닝:\n  EfficientNet-B3 결합 및 단계적 Backbone Unfreezing\n  ➔ 최종 Validation Loss 0.16 달성"
]
for b in bullets_l:
    p = tf_l.add_paragraph()
    p.text = b
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_SUB
    p.space_before = Pt(10)

create_card(s5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9))
tb_r = s5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
p = tf_r.paragraphs[0]
p.text = "⚙️ 72회 추론 그리드 서치 (Grid Search)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

bullets_r = [
    "• 탐색 파라미터 조합 (총 72회 실험):\n  - 해상도: 640 / 768\n  - Bounding Box 패딩 비율: 0.08 ~ 0.12\n  - Detector / Classifier 가중치 지수 (Beta)",
    "• 정량적 성능 향상 비교:\n  - 초기 베이스라인 mAP: 0.742\n  - 최종 최적화 mAP50-95: 0.877 (약 18.2% 성능 향상)\n  - mAP50: 0.887 / Kaggle Public Score: 0.987 달성"
]
for b in bullets_r:
    p = tf_r.add_paragraph()
    p.text = b
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_SUB
    p.space_before = Pt(10)

# ==========================================
# [Slide 6] 프로젝트 2 아키텍처 비교 (수치 명시)
# ==========================================
s6 = prs.slides.add_slide(blank_layout)
apply_background(s6)
add_header(s6, "Project 02 • Comparison & Insight", "아키텍처별 성능 비교 및 핵심 인사이트")

rows, cols = 4, 4
table_shape6 = s6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.6))
t6 = table_shape6.table
t6.columns[0].width = Inches(2.8)
t6.columns[1].width = Inches(3.2)
t6.columns[2].width = Inches(3.0)
t6.columns[3].width = Inches(2.73)

headers6 = ["모델 아키텍처", "장점 및 특징", "한계점 (Bottleneck)", "검증 지표 (성능 비교)"]
for col_idx, h_text in enumerate(headers6):
    cell = t6.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER_BG
    p = cell.text_frame.paragraphs[0]
    p.text = h_text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

p2_table_data = [
    ("YOLO26m (1-Stage 단독)", "빠른 추론 속도", "소수 클래스 간 오분류 및\n미세 각인 인식률 저하", "mAP 0.742 (초기 베이스라인)"),
    ("Faster R-CNN (단독)", "우수한 BBox 검출률", "다중 클래스 동시 수렴 지연\n파라미터 연산 과다", "mAP 0.812 (Val Loss 0.16)"),
    ("2-Stage 하이브리드\n(YOLO + EfficientNet)", "위치 검출과 정밀 분류 책임 분리\n불균형 가중 샘플러 적용", "2단계 파이프라인 연산 필요", "COCO mAP 0.877 (18.2% ↑)\nKaggle 0.987 최고점 달성")
]

for row_idx, row_data in enumerate(p2_table_data, start=1):
    for col_idx, text in enumerate(row_data):
        cell = t6.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_SUB
        if col_idx == 3 and row_idx == 3:
            p.font.bold = True
            p.font.color.rgb = ACCENT_GREEN

create_card(s6, Inches(0.8), Inches(5.7), Inches(11.73), Inches(1.1), bg_color=ACCENT_BLUE, border_color=PRIMARY_BLUE)
tb_sub6 = s6.shapes.add_textbox(Inches(1.1), Inches(5.8), Inches(11.1), Inches(0.8))
tf_sub6 = tb_sub6.text_frame
tf_sub6.word_wrap = True
p = tf_sub6.paragraphs[0]
p.text = "💡 핵심 결론: 단일 모델 한계를 극복하는 2-Stage 분리 아키텍처로 초기 베이스라인 대비 mAP 18.2% 향상 및 Kaggle 0.987을 견인함"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

# ==========================================
# [Slide 7] 엔지니어링 역량 요약
# ==========================================
s7 = prs.slides.add_slide(blank_layout)
apply_background(s7)
add_header(s7, "Summary & Vision", "엔지니어링 핵심 역량 및 입사 후 포부")

cards = [
    ("지표 중심의 모델링", "• 72회 하이퍼파라미터 그리드 서치\n• 베이스라인 대비 mAP 18.2% 향상\n• 불균형 교정으로 정확도 14.5% 개선"),
    ("안정적인 서빙 최적화", "• FastAPI 비동기 BackgroundTasks\n• 체감 응답 지연 60.5% 단축 달성\n• Lifespan 인메모리 Cold Start 제거"),
    ("MLOps 관측 & 비용 통제", "• Langfuse @observe 실시간 트레이싱\n• 건당 $0.0082 비용 통제 체계\n• 지속 가능한 서비스 파이프라인 구축")
]
for i, (ctitle, cdesc) in enumerate(cards):
    left = Inches(0.8 + i * 4.0)
    create_card(s7, left, Inches(1.8), Inches(3.73), Inches(3.4))
    tb_c = s7.shapes.add_textbox(left + Inches(0.2), Inches(2.1), Inches(3.33), Inches(2.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    p1 = tf_c.paragraphs[0]
    p1.text = ctitle
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_BLUE
    p2 = tf_c.add_paragraph()
    p2.text = cdesc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SUB
    p2.space_before = Pt(10)

create_card(s7, Inches(0.8), Inches(5.5), Inches(11.73), Inches(1.3), bg_color=ACCENT_BLUE, border_color=PRIMARY_BLUE)
tb_goal = s7.shapes.add_textbox(Inches(1.1), Inches(5.65), Inches(11.1), Inches(1.0))
tf_goal = tb_goal.text_frame
tf_goal.word_wrap = True
p = tf_goal.paragraphs[0]
p.text = "🎯 입사 후 포부"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

p = tf_goal.add_paragraph()
p.text = "단순 모델 개발에 그치지 않고 배포 환경의 응답성과 비용 효율까지 책임지는 엔지니어로서, 실제 비즈니스 환경에서 안정적으로 동작하는 AI 솔루션을 구축하겠습니다."
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = TEXT_DARK
p.space_before = Pt(3)

# 파일 저장
output_path = "AI_엔지니어_조용준_포트폴리오.pptx"
prs.save(output_path)
print(f"✅ 성공적으로 생성되었습니다: {output_path}")